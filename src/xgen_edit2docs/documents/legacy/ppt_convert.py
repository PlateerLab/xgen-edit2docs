"""PowerPoint 97 바이너리(.ppt) → PPTX 변환.

'PowerPoint Document' 스트림은 8바이트 헤더(recVer/instance 2B, recType 2B,
recLen 4B) 레코드 트리다. 슬라이드 텍스트의 정본은 DocumentContainer(1000)
안의 SlideListWithText(4080) — SlidePersistAtom(1011)이 슬라이드 경계를,
TextHeaderAtom(3999)이 자리 종류(0/5=제목, 1/6=본문 …)를, TextCharsAtom
(4000, UTF-16LE)/TextBytesAtom(4008, 8비트=UTF-16 하위바이트)이 내용을,
**StyleTextPropAtom(4001)** 이 직전 텍스트 아톰의 문단/문자 서식 런을 준다.

StyleTextPropAtom 레이아웃 ([MS-PPT] 2.9.20/2.9.44, Apache POI HSLF
TextPropCollection 대조 — reference_data/poi):

    문단 런*  { textLen(4) indentLevel(2) mask(4) props(mask 순서별) }
    문자 런*  { textLen(4) mask(4) props }
    문자 props 순서: charFlags(mask&0xFFFF, 2B — bit0 bold/bit1 italic/
      bit2 underline/bit8 strike) → font.index(0x10000,2) → asian(0x200000,2)
      → ansi(0x400000,2) → symbol(0x800000,2) → font.size(0x20000,2)
      → font.color(0x40000,4 — RGB, 상위 바이트 0xFE/0xFF 는 스킴색)
      → superscript(0x80000,2)
    문단 props 순서: paraFlags(0xF,2) → bullet.char(0x80,2) →
      bullet.font(0x10,2) → bullet.size(0x40,2) → bullet.color(0x20,4) →
      alignment(0x800,2 — 0 좌/1 중앙/2 우/3 양쪽) → linespacing(0x1000,2)
      → spacebefore(0x2000,2) → spaceafter(0x4000,2) → text.offset(0x100,2)
      → bullet.offset(0x400,2) → defaultTabSize(0x8000,2) →
      tabstops(0x100000, 2+n×4 가변) → fontAlign(0x10000,2) →
      wrapFlags(0xE0000,2) → textDirection(0x200000,2)

슬라이드 크기는 DocumentAtom(1001)의 master unit(1/576 inch) 값.

충실도 범위: 슬라이드 수·순서, 자리별 텍스트(제목/본문), 문단 정렬·
들여쓰기 수준, 문자 런 스타일(굵게/기울임/밑줄/취소선/크기/색).
도형 좌표(Escher)·이미지·표는 범위 밖 — 텍스트 유실은 없다.
"""

from __future__ import annotations

import io
import struct
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

from . import LegacyConvertError

_RT_DOCUMENT = 1000
_RT_DOCUMENT_ATOM = 1001
_RT_SLIDE_PERSIST_ATOM = 1011
_RT_SLIDE_LIST_WITH_TEXT = 4080
_RT_TEXT_HEADER_ATOM = 3999
_RT_TEXT_CHARS_ATOM = 4000
_RT_STYLE_TEXT_PROP_ATOM = 4001
_RT_TEXT_BYTES_ATOM = 4008

_TITLE_TYPES = {0, 5}  # title / center title

_MASTER_PER_INCH = 576.0
_EMU_PER_INCH = 914400.0


def _iter_records(data: bytes, start: int = 0, end: int | None = None):
    """(rec_ver, rec_type, payload_start, payload_len) 평면 순회."""
    pos = start
    n = len(data) if end is None else end
    while pos + 8 <= n:
        ver_inst, rec_type, rec_len = struct.unpack_from("<HHI", data, pos)
        payload_start = pos + 8
        if rec_len > n - payload_start:
            return
        yield (ver_inst & 0x000F), rec_type, payload_start, rec_len
        pos = payload_start + rec_len


# ── StyleTextPropAtom 해석 ─────────────────────────────────────


@dataclass
class _CharStyle:
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    strike: Optional[bool] = None
    size_pt: Optional[float] = None
    color: Optional[str] = None  # RRGGBB


@dataclass
class _ParaStyle:
    indent: int = 0
    align: Optional[int] = None  # 0 left / 1 center / 2 right / 3 justify


#: (mask, size) — 문자 props, POI characterTextPropTypes 순서.
_CHAR_PROPS: List[Tuple[int, int]] = [
    (0xFFFF, 2),      # charFlags
    (0x10000, 2),     # font.index
    (0x200000, 2),    # asian.font.index
    (0x400000, 2),    # ansi.font.index
    (0x800000, 2),    # symbol.font.index
    (0x20000, 2),     # font.size
    (0x40000, 4),     # font.color
    (0x80000, 2),     # superscript
]

#: (mask, size) — 문단 props (tabstops 는 가변이라 특별 취급).
_PARA_PROPS: List[Tuple[int, int]] = [
    (0xF, 2),         # paraFlags
    (0x80, 2),        # bullet.char
    (0x10, 2),        # bullet.font
    (0x40, 2),        # bullet.size
    (0x20, 4),        # bullet.color
    (0x800, 2),       # alignment
    (0x1000, 2),      # linespacing
    (0x2000, 2),      # spacebefore
    (0x4000, 2),      # spaceafter
    (0x100, 2),       # text.offset
    (0x400, 2),       # bullet.offset
    (0x8000, 2),      # defaultTabSize
    (0x100000, -1),   # tabstops — 2 + n×4
    (0x10000, 2),     # fontAlign
    (0xE0000, 2),     # wrapFlags
    (0x200000, 2),    # textDirection
]


def _parse_style_atom(payload: bytes, text_len: int
                      ) -> Tuple[List[Tuple[int, _ParaStyle]],
                                 List[Tuple[int, _CharStyle]]]:
    """→ ([(문단 런 길이, 스타일)], [(문자 런 길이, 스타일)]).

    깨진/모르는 마스크를 만나면 그 지점까지 해석한 결과만 돌려준다 —
    스타일은 최선 노력, 텍스트는 이미 안전하다.
    """
    pos, n = 0, len(payload)
    para_runs: List[Tuple[int, _ParaStyle]] = []
    char_runs: List[Tuple[int, _CharStyle]] = []
    try:
        handled = 0
        while pos + 10 <= n and handled <= text_len:
            (run_len,) = struct.unpack_from("<I", payload, pos)
            (indent,) = struct.unpack_from("<H", payload, pos + 4)
            (mask,) = struct.unpack_from("<I", payload, pos + 6)
            pos += 10
            st = _ParaStyle(indent=min(max(indent, 0), 8))
            for m, size in _PARA_PROPS:
                if not (mask & m):
                    continue
                if size < 0:  # tabstops
                    if pos + 2 > n:
                        raise ValueError
                    (cnt,) = struct.unpack_from("<H", payload, pos)
                    size = 2 + cnt * 4
                if pos + size > n:
                    raise ValueError
                if m == 0x800:
                    (st.align,) = struct.unpack_from("<H", payload, pos)
                pos += size
            para_runs.append((run_len, st))
            handled += run_len
            if handled >= text_len + 1:
                break

        handled = 0
        while pos + 8 <= n and handled <= text_len:
            (run_len,) = struct.unpack_from("<I", payload, pos)
            (mask,) = struct.unpack_from("<I", payload, pos + 4)
            pos += 8
            st = _CharStyle()
            for m, size in _CHAR_PROPS:
                if not (mask & m):
                    continue
                if pos + size > n:
                    raise ValueError
                if m == 0xFFFF:
                    (flags,) = struct.unpack_from("<H", payload, pos)
                    # mask 비트가 선언된 플래그만 유효 (POI setValueWithMask)
                    if mask & 0x0001:
                        st.bold = bool(flags & 0x0001)
                    if mask & 0x0002:
                        st.italic = bool(flags & 0x0002)
                    if mask & 0x0004:
                        st.underline = bool(flags & 0x0004)
                    if mask & 0x0100:
                        st.strike = bool(flags & 0x0100)
                elif m == 0x20000:
                    (sz,) = struct.unpack_from("<H", payload, pos)
                    if 1 <= sz <= 999:
                        st.size_pt = float(sz)
                elif m == 0x40000:
                    r, g, b, kind = struct.unpack_from("<4B", payload, pos)
                    if kind not in (0xFE, 0xFF):  # 스킴색 참조는 범위 밖
                        st.color = f"{r:02X}{g:02X}{b:02X}"
                pos += size
            char_runs.append((run_len, st))
            handled += run_len
            if handled >= text_len + 1:
                break
    except (ValueError, struct.error):
        pass
    return para_runs, char_runs


# ── 슬라이드 수집 ──────────────────────────────────────────────


@dataclass
class _TextBlock:
    is_title: bool
    text: str
    para_runs: List[Tuple[int, _ParaStyle]] = field(default_factory=list)
    char_runs: List[Tuple[int, _CharStyle]] = field(default_factory=list)


@dataclass
class _SlideText:
    blocks: List[_TextBlock] = field(default_factory=list)


def _decode_bytes_atom(payload: bytes) -> str:
    # TextBytesAtom — 각 바이트가 UTF-16 코드포인트의 하위 바이트다.
    return payload.decode("latin-1")


def _collect_slides(data: bytes) -> tuple[List[_SlideText], tuple[int, int]]:
    slide_size_mu = (9144, 6858)  # 10in × 7.5in 기본
    slides: List[_SlideText] = []

    def walk(start: int, end: int, in_sltwt: bool) -> None:
        nonlocal slide_size_mu
        cur_type = 1  # TextHeader 없이 오는 텍스트는 본문 취급
        for rec_ver, rec_type, p_start, p_len in _iter_records(data, start, end):
            if rec_ver == 0xF:  # 컨테이너 — 재귀
                walk(p_start, p_start + p_len,
                     in_sltwt or rec_type == _RT_SLIDE_LIST_WITH_TEXT)
                continue
            if rec_type == _RT_DOCUMENT_ATOM and p_len >= 8:
                w, h = struct.unpack_from("<ii", data, p_start)
                if 576 <= w <= 576 * 100 and 576 <= h <= 576 * 100:
                    slide_size_mu = (w, h)
            if not in_sltwt:
                continue
            if rec_type == _RT_SLIDE_PERSIST_ATOM:
                slides.append(_SlideText())
                cur_type = 1
            elif rec_type == _RT_TEXT_HEADER_ATOM and p_len >= 4:
                (cur_type,) = struct.unpack_from("<I", data, p_start)
            elif rec_type in (_RT_TEXT_CHARS_ATOM, _RT_TEXT_BYTES_ATOM):
                if not slides:
                    slides.append(_SlideText())
                raw = data[p_start:p_start + p_len]
                text = (raw.decode("utf-16le", errors="replace")
                        if rec_type == _RT_TEXT_CHARS_ATOM
                        else _decode_bytes_atom(raw))
                slides[-1].blocks.append(
                    _TextBlock(is_title=cur_type in _TITLE_TYPES, text=text))
            elif rec_type == _RT_STYLE_TEXT_PROP_ATOM:
                # 직전 텍스트 아톰의 서식 런
                if slides and slides[-1].blocks:
                    block = slides[-1].blocks[-1]
                    block.para_runs, block.char_runs = _parse_style_atom(
                        data[p_start:p_start + p_len], len(block.text))

    walk(0, len(data), False)
    return slides, slide_size_mu


# ── PPTX 조립 ──────────────────────────────────────────────────


def _char_style_at(block: _TextBlock, pos: int) -> _CharStyle:
    acc = 0
    for run_len, st in block.char_runs:
        if pos < acc + run_len:
            return st
        acc += run_len
    return _CharStyle()


def _char_boundaries(block: _TextBlock) -> List[int]:
    out, acc = [], 0
    for run_len, _st in block.char_runs:
        acc += run_len
        out.append(acc)
    return out


def _para_style_at(block: _TextBlock, pos: int) -> _ParaStyle:
    acc = 0
    for run_len, st in block.para_runs:
        if pos < acc + run_len:
            return st
        acc += run_len
    return _ParaStyle()


def ppt_to_pptx(content: bytes) -> bytes:
    import olefile
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt

    if not olefile.isOleFile(io.BytesIO(content)):
        raise LegacyConvertError("ppt 가 아닙니다 (OLE 복합문서 아님)")
    ole = olefile.OleFileIO(io.BytesIO(content))
    try:
        if not ole.exists("PowerPoint Document"):
            raise LegacyConvertError("ppt 'PowerPoint Document' 스트림이 없습니다")
        data = ole.openstream("PowerPoint Document").read()
    finally:
        ole.close()

    slides, (w_mu, h_mu) = _collect_slides(data)
    if not slides:
        raise LegacyConvertError("ppt 에서 슬라이드 텍스트를 찾지 못했습니다")

    prs = Presentation()
    prs.slide_width = Emu(int(w_mu * _EMU_PER_INCH / _MASTER_PER_INCH))
    prs.slide_height = Emu(int(h_mu * _EMU_PER_INCH / _MASTER_PER_INCH))
    blank = prs.slide_layouts[6]

    margin = Emu(int(0.5 * _EMU_PER_INCH))
    title_h = Emu(int(1.1 * _EMU_PER_INCH))
    content_w = prs.slide_width - margin * 2

    _ALIGN = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER,
              2: PP_ALIGN.RIGHT, 3: PP_ALIGN.JUSTIFY}

    def emit_block(tf, block: _TextBlock, default_pt: float,
                   default_bold: bool, first_para_used: bool) -> bool:
        """텍스트 블록 → 문단들(\\r 경계), 런 스타일/정렬/레벨 반영."""
        boundaries = _char_boundaries(block)
        pos = 0
        for para_text in block.text.split("\r"):
            para = tf.paragraphs[0] if not first_para_used else tf.add_paragraph()
            first_para_used = True
            pst = _para_style_at(block, pos)
            if pst.align in _ALIGN:
                para.alignment = _ALIGN[pst.align]
            if pst.indent:
                para.level = min(pst.indent, 4)
            # 문자 런 경계 + 줄바꿈(0x0B) 지점으로 조각 낸다
            cuts = sorted({pos, pos + len(para_text)} | {
                b for b in boundaries if pos < b < pos + len(para_text)} | {
                pos + i for i, ch in enumerate(para_text) if ch == "\x0b"} | {
                pos + i + 1 for i, ch in enumerate(para_text) if ch == "\x0b"})
            for a, b in zip(cuts, cuts[1:]):
                piece = block.text[a:b]
                if piece == "\x0b":
                    para.add_line_break()
                    continue
                if not piece:
                    continue
                st = _char_style_at(block, a)
                run = para.add_run()
                run.text = piece.replace("\x0b", "\n")
                run.font.size = Pt(st.size_pt if st.size_pt else default_pt)
                run.font.bold = st.bold if st.bold is not None else default_bold
                if st.italic is not None:
                    run.font.italic = st.italic
                if st.underline is not None:
                    run.font.underline = st.underline
                if st.strike:
                    # python-pptx 에 strike API 가 없다 — rPr 속성 직접
                    run.font._rPr.set("strike", "sngStrike")
                if st.color:
                    run.font.color.rgb = RGBColor.from_string(st.color)
            pos += len(para_text) + 1  # + \r
        return first_para_used

    for st_slide in slides:
        slide = prs.slides.add_slide(blank)
        y = margin
        titles = [b for b in st_slide.blocks if b.is_title]
        bodies = [b for b in st_slide.blocks if not b.is_title]
        if titles:
            box = slide.shapes.add_textbox(margin, y, content_w, title_h)
            tf = box.text_frame
            tf.word_wrap = True
            used = False
            for block in titles:
                used = emit_block(tf, block, 28.0, True, used)
            y = y + title_h
        if bodies:
            body_h = prs.slide_height - y - margin
            box = slide.shapes.add_textbox(margin, y, content_w, body_h)
            tf = box.text_frame
            tf.word_wrap = True
            used = False
            for block in bodies:
                used = emit_block(tf, block, 16.0, False, used)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
