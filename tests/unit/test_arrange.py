"""arrange_doc — deterministic structural edits (slides / sheets).

Acceptance is a reopen with the host Office library and, for the agent
surface, dispatch through run_tool.
"""

from __future__ import annotations

import io


def _deck() -> bytes:
    import os
    import tempfile

    from xgen_edit2docs import build_doc

    d = tempfile.mkdtemp()
    p = os.path.join(d, "deck.pptx")
    build_doc(
        {"slides": [{"title": "S0"}, {"title": "S1"}, {"title": "S2"}]}, p
    )
    return open(p, "rb").read()


def _titles(data: bytes) -> list[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out = []
    for s in prs.slides:
        out.append(
            next(
                (
                    sh.text_frame.text.split("\n")[0]
                    for sh in s.shapes
                    if sh.has_text_frame and sh.text_frame.text.strip()
                ),
                "",
            )
        )
    return out


def _wb() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.title = "Sales"
    wb.active["A1"] = 1
    wb.create_sheet("Notes")["A1"] = "hi"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestPptxArrange:
    def test_duplicate(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        res = arrange_doc(str(p), [{"op": "duplicate", "target": 0, "to": 3}])
        assert res.applied == 1
        assert _titles(res.path.read_bytes()) == ["S0", "S1", "S2", "S0"]

    def test_move(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        res = arrange_doc(str(p), [{"op": "move", "target": 2, "to": 0}])
        assert _titles(res.path.read_bytes()) == ["S2", "S0", "S1"]

    def test_delete(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        res = arrange_doc(str(p), [{"op": "delete", "target": 1}])
        assert _titles(res.path.read_bytes()) == ["S0", "S2"]

    def test_sequential_ops_resolve_against_current_state(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        # delete S1 (→ [S0,S2]), then move index1 (S2) to front
        res = arrange_doc(
            str(p),
            [{"op": "delete", "target": 1}, {"op": "move", "target": 1, "to": 0}],
        )
        assert res.applied == 2
        assert _titles(res.path.read_bytes()) == ["S2", "S0"]

    def test_not_found_and_invalid_are_best_effort(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        res = arrange_doc(
            str(p),
            [
                {"op": "delete", "target": 99},  # not_found
                {"op": "frobnicate", "target": 0},  # invalid
                {"op": "move", "target": 0},  # invalid: no 'to'
                {"op": "duplicate", "target": 0},  # applied
            ],
        )
        statuses = [r["status"] for r in res.results]
        assert statuses == ["not_found", "invalid", "invalid", "applied"]
        assert res.applied == 1

    def test_no_op_applied_leaves_original(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        res = arrange_doc(str(p), [{"op": "delete", "target": 99}])
        assert res.applied == 0
        assert res.path == p  # points back at the untouched input


class TestXlsxArrange:
    def test_duplicate_move_rename_delete(self, tmp_path):
        from openpyxl import load_workbook

        from xgen_edit2docs import arrange_doc

        p = tmp_path / "b.xlsx"
        p.write_bytes(_wb())
        res = arrange_doc(
            str(p),
            [
                {"op": "duplicate", "target": "Sales", "name": "Sales copy", "to": 1},
                {"op": "rename", "target": "Notes", "name": "Memos"},
                {"op": "move", "target": "Memos", "to": 0},
            ],
        )
        assert res.applied == 3
        names = load_workbook(io.BytesIO(res.path.read_bytes())).sheetnames
        assert names == ["Memos", "Sales", "Sales copy"]

    def test_delete_last_sheet_refused(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "b.xlsx"
        p.write_bytes(_wb())
        res = arrange_doc(
            str(p),
            [{"op": "delete", "target": "Notes"}, {"op": "delete", "target": "Sales"}],
        )
        assert [r["status"] for r in res.results] == ["applied", "refused"]

    def test_duplicate_without_name_is_invalid(self, tmp_path):
        from xgen_edit2docs import arrange_doc

        p = tmp_path / "b.xlsx"
        p.write_bytes(_wb())
        res = arrange_doc(str(p), [{"op": "duplicate", "target": "Sales"}])
        assert res.results[0]["status"] == "invalid"


class TestDocxRejected:
    def test_docx_has_no_structure(self, tmp_path):
        from xgen_edit2docs import arrange_doc
        from xgen_edit2docs.documents.docx_engine import docx_from_markdown

        p = tmp_path / "d.docx"
        p.write_bytes(docx_from_markdown("# Hi"))
        res = arrange_doc(str(p), [{"op": "move", "target": 0, "to": 1}])
        assert res.results[0]["status"] == "invalid"


class TestAgentSurface:
    def test_run_tool_arrange(self, tmp_path):
        from xgen_edit2docs.agent_tools import TOOL_NAMES, run_tool

        assert "arrange_doc" in TOOL_NAMES
        p = tmp_path / "d.pptx"
        p.write_bytes(_deck())
        out = run_tool(
            "arrange_doc",
            {
                "doc": str(p),
                "ops": [{"op": "duplicate", "target": 0}],
                "output": str(tmp_path / "o.pptx"),
            },
        )
        assert out["applied"] == 1
        assert "warnings" in out

    def test_arrange_guide_topic_resolves(self):
        from xgen_edit2docs.agent_guide import doc_guide

        g = doc_guide("arrange")
        assert "arrange_doc" in g["guide"]
        assert "arrange" in doc_guide()["topics"]
