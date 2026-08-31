"""documents/legacy — hwp/hwpx/doc/xls/ppt → OOXML → 페이지 SVG E2E.

픽스처는 전부 손으로 조립한다 — 실제 포맷 스펙(레코드/스트림 구조)을
그대로 따라 만들므로, 파서가 스펙의 어느 지점을 읽는지가 테스트에
드러난다. OLE 복합문서(CFB)는 아래 미니 라이터로 만든다 (모든 스트림을
4096B 이상으로 패딩해 ministream 없이 정규 섹터 체인만 쓴다 — 각 파서는
트레일링 패딩을 관용한다는 것 자체가 계약이다).
"""

from __future__ import annotations

import io
import struct
import zipfile
import zlib
from pathlib import Path

import pytest

from xgen_edit2docs.documents.legacy import LegacyConvertError, convert_to_ooxml
from xgen_edit2docs.simple import render_doc

# ─────────────────────────────────────────────────────────────
# 미니 CFB(v3) 라이터 — 정규 섹터 체인 전용
# ─────────────────────────────────────────────────────────────

_ENDOFCHAIN = 0xFFFFFFFE
_FATSECT = 0xFFFFFFFD
_FREESECT = 0xFFFFFFFF
_NOSTREAM = 0xFFFFFFFF
_SECT = 512


def build_cfb(streams: dict[str, bytes]) -> bytes:
    """{이름: 내용} → 최소 CFB. 스트림은 4096B 로 패딩(ministream 회피)."""
    padded = {name: data + b"\x00" * max(0, 4096 - len(data))
              for name, data in streams.items()}

    # 디렉터리: 루트 + 스토리지('A/B' 의 A) + 스트림 엔트리
    n_storages = len({n.split("/", 1)[0] for n in padded if "/" in n})
    n_entries = 1 + n_storages + len(padded)
    n_dir_sectors = (n_entries * 128 + _SECT - 1) // _SECT

    chunks: list[bytes] = []          # 데이터 섹터 (FAT 뒤에 온다)
    chains: list[tuple[int, int]] = []  # (start, n_sectors) per stream

    # 섹터 배치: [FAT…][DIR…][stream data…]
    # FAT 섹터 수는 자기 자신 포함 총 섹터로부터 고정점 계산.
    stream_sectors = sum((len(d) + _SECT - 1) // _SECT for d in padded.values())
    total_data = n_dir_sectors + stream_sectors
    n_fat = 1
    while n_fat * 128 < n_fat + total_data:
        n_fat += 1

    dir_start = n_fat
    stream_start = dir_start + n_dir_sectors

    fat: list[int] = []
    for i in range(n_fat):
        fat.append(_FATSECT)
    for i in range(n_dir_sectors):
        fat.append(dir_start + i + 1 if i + 1 < n_dir_sectors else _ENDOFCHAIN)

    cursor = stream_start
    entries: list[bytes] = []

    def dir_entry(name: str, etype: int, start: int, size: int,
                  child: int = _NOSTREAM, right: int = _NOSTREAM) -> bytes:
        raw_name = name.encode("utf-16le")[:62]
        e = bytearray(128)
        e[0:len(raw_name)] = raw_name
        struct.pack_into("<H", e, 64, len(raw_name) + 2)
        e[66] = etype        # 1=storage, 2=stream, 5=root
        e[67] = 1            # black
        struct.pack_into("<I", e, 68, _NOSTREAM)  # left
        struct.pack_into("<I", e, 72, right)
        struct.pack_into("<I", e, 76, child)
        struct.pack_into("<I", e, 116, start)
        struct.pack_into("<Q", e, 120, size)
        return bytes(e)

    names = list(padded.keys())
    for idx, name in enumerate(names):
        data = padded[name]
        n_sec = (len(data) + _SECT - 1) // _SECT
        for s in range(n_sec):
            fat.append(cursor + s + 1 if s + 1 < n_sec else _ENDOFCHAIN)
        chains.append((cursor, n_sec))
        chunks.append(data + b"\x00" * (n_sec * _SECT - len(data)))
        cursor += n_sec

    # 디렉터리 트리 — 'A/B' 이름은 스토리지 A 아래 스트림 B 다 (HWP 의
    # BodyText/Section0 처럼). 루트 child → 최상위 체인, 스토리지 child →
    # 하위 스트림 체인 (형제는 right 포인터로 잇는다 — 퇴화 트리지만 유효).
    top_level: list[tuple[str, list[int]]] = []  # (이름, [stream idx]) — 스트림이면 idx 1개
    storages: dict[str, list[int]] = {}
    order: list[str] = []
    for idx, name in enumerate(names):
        if "/" in name:
            stg, _leaf = name.split("/", 1)
            if stg not in storages:
                storages[stg] = []
                order.append("STG:" + stg)
            storages[stg].append(idx)
        else:
            order.append("STR:" + str(idx))

    # 엔트리 배열을 먼저 자리만 잡고(인덱스 확정) 포인터를 나중에 채운다.
    specs: list[dict] = [dict(name="Root Entry", etype=5, start=_ENDOFCHAIN,
                              size=0, child=_NOSTREAM, right=_NOSTREAM)]
    entry_of: dict[str, int] = {}
    for key in order:
        if key.startswith("STG:"):
            stg = key[4:]
            entry_of[key] = len(specs)
            specs.append(dict(name=stg, etype=1, start=_ENDOFCHAIN, size=0,
                              child=_NOSTREAM, right=_NOSTREAM))
            for sidx in storages[stg]:
                leaf = names[sidx].split("/", 1)[1]
                entry_of[f"LEAF:{sidx}"] = len(specs)
                start, _n = chains[sidx]
                specs.append(dict(name=leaf, etype=2, start=start,
                                  size=len(padded[names[sidx]]),
                                  child=_NOSTREAM, right=_NOSTREAM))
        else:
            sidx = int(key[4:])
            entry_of[key] = len(specs)
            start, _n = chains[sidx]
            specs.append(dict(name=names[sidx], etype=2, start=start,
                              size=len(padded[names[sidx]]),
                              child=_NOSTREAM, right=_NOSTREAM))

    # 포인터 배선
    top_ids = [entry_of[k] for k in order]
    if top_ids:
        specs[0]["child"] = top_ids[0]
        for a, b in zip(top_ids, top_ids[1:]):
            specs[a]["right"] = b
    for key in order:
        if key.startswith("STG:"):
            stg = key[4:]
            kid_ids = [entry_of[f"LEAF:{i}"] for i in storages[stg]]
            specs[entry_of[key]]["child"] = kid_ids[0]
            for a, b in zip(kid_ids, kid_ids[1:]):
                specs[a]["right"] = b

    for s in specs:
        entries.append(dir_entry(s["name"], s["etype"], s["start"], s["size"],
                                 child=s["child"], right=s["right"]))

    dir_blob = b"".join(entries)
    dir_blob += b"\x00" * (n_dir_sectors * _SECT - len(dir_blob))

    fat += [_FREESECT] * (n_fat * 128 - len(fat))
    fat_blob = struct.pack(f"<{len(fat)}I", *fat)

    header = bytearray(_SECT)
    header[0:8] = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    struct.pack_into("<H", header, 24, 0x003E)   # minor
    struct.pack_into("<H", header, 26, 0x0003)   # major (v3)
    struct.pack_into("<H", header, 28, 0xFFFE)   # byte order
    struct.pack_into("<H", header, 30, 9)        # sector shift
    struct.pack_into("<H", header, 32, 6)        # mini shift
    struct.pack_into("<I", header, 44, n_fat)
    struct.pack_into("<I", header, 48, dir_start)
    struct.pack_into("<I", header, 56, 4096)     # mini cutoff
    struct.pack_into("<I", header, 60, _ENDOFCHAIN)  # first minifat
    struct.pack_into("<I", header, 64, 0)
    struct.pack_into("<I", header, 68, _ENDOFCHAIN)  # first difat
    struct.pack_into("<I", header, 72, 0)
    for i in range(109):
        struct.pack_into("<I", header, 76 + i * 4,
                         i if i < n_fat else _FREESECT)

    return bytes(header) + fat_blob + dir_blob + b"".join(chunks)


# ─────────────────────────────────────────────────────────────
# 포맷별 픽스처 조립
# ─────────────────────────────────────────────────────────────


def _rec(tagid: int, level: int, payload: bytes) -> bytes:
    hdr = (tagid & 0x3FF) | ((level & 0x3FF) << 10) | ((len(payload) & 0xFFF) << 20)
    return struct.pack("<I", hdr) + payload


def _utf16(s: str) -> bytes:
    return s.encode("utf-16le")


def make_hwp(text: str = "한글 렌더 검증", bold_text: str = "굵은 부분") -> bytes:
    header = bytearray(256)
    header[0:32] = b"HWP Document File".ljust(32, b"\x00")
    struct.pack_into("<I", header, 32, 0x05000000)  # version
    struct.pack_into("<I", header, 36, 0x1)         # compressed

    # DocInfo: char shape 2개 — 보통(10pt), 굵게(14pt, 빨강)
    def char_shape(size_pt: float, flags: int, color_bgr: int) -> bytes:
        p = bytearray(72)
        struct.pack_into("<i", p, 42, int(size_pt * 100))
        struct.pack_into("<I", p, 46, flags)
        struct.pack_into("<I", p, 52, color_bgr)
        return _rec(0x15, 0, bytes(p))

    docinfo = char_shape(10.0, 0x0, 0x000000) + char_shape(14.0, 0x2, 0x0000FF)

    # Section0: PAGE_DEF + 문단 2개 (두 번째는 charshape 1 굵게) + 1×2 표
    page = struct.pack("<6I", 59528, 84188, 8504, 8504, 5668, 4252) + b"\x00" * 16
    para1_text = _utf16(text)
    para2_text = _utf16(bold_text)
    tbl_flags = struct.pack("<I", 0) + struct.pack("<HH", 1, 2) + b"\x00" * 8
    cell_a = _utf16("셀A")
    cell_b = _utf16("셀B")
    section = b"".join([
        _rec(0x49, 1, page),
        _rec(0x42, 0, b"\x00" * 16),               # PARA_HEADER
        _rec(0x43, 1, para1_text),                 # PARA_TEXT
        _rec(0x44, 1, struct.pack("<II", 0, 0)),   # PARA_CHAR_SHAPE → shape 0
        _rec(0x42, 0, b"\x00" * 16),
        _rec(0x43, 1, para2_text),
        _rec(0x44, 1, struct.pack("<II", 0, 1)),   # shape 1 (bold/red)
        # 표: CTRL_HEADER('tbl ') → TABLE(1×2) → [LIST_HEADER + 문단]×2
        _rec(0x47, 0, b" lbt" + b"\x00" * 8),
        _rec(0x4D, 1, tbl_flags),
        _rec(0x48, 1, b"\x00" * 8),
        _rec(0x42, 2, b"\x00" * 16),
        _rec(0x43, 3, cell_a),
        _rec(0x48, 1, b"\x00" * 8),
        _rec(0x42, 2, b"\x00" * 16),
        _rec(0x43, 3, cell_b),
    ])

    def comp(b: bytes) -> bytes:
        co = zlib.compressobj(6, zlib.DEFLATED, -15)
        return co.compress(b) + co.flush()

    return build_cfb({
        "FileHeader": bytes(header),
        "DocInfo": comp(docinfo),
        "BodyText/Section0": comp(section),
    })


def _tiny_png() -> bytes:
    """1×1 RGB PNG — CRC 까지 규격대로 조립한다."""
    def chunk(kind: bytes, body: bytes) -> bytes:
        return (struct.pack(">I", len(body)) + kind + body
                + struct.pack(">I", zlib.crc32(kind + body) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    idat = zlib.compress(b"\x00\xff\x00\x00")  # filter 0 + RGB(255,0,0)
    return (b"\x89PNG\r\n\x1a\x0a" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", idat) + chunk(b"IEND", b""))


def make_hwp_rich() -> bytes:
    """심화 픽스처 — 병합/배경/정렬/취소선/글꼴/이미지/머리말을 전부 싣는다.

    실파일 규격 그대로: 컨트롤(CTRL_HEADER)은 앵커 문단의 자식 레벨,
    셀 LIST_HEADER 는 표 75 셀 속성(col/row/colspan/rowspan/크기/borderfill)
    을 온전히 갖는다.
    """
    header = bytearray(256)
    header[0:32] = b"HWP Document File".ljust(32, b"\x00")
    struct.pack_into("<I", header, 32, 0x05000000)
    struct.pack_into("<I", header, 36, 0x1)  # compressed

    def bstr(s: str) -> bytes:
        return struct.pack("<H", len(s)) + s.encode("utf-16le")

    # ── DocInfo ──
    id_mappings = struct.pack("<8i", 1, 1, 0, 0, 0, 0, 0, 0)  # binData 1, ko 글꼴 1
    bin_data = struct.pack("<HH", 0x1, 1) + bstr("png")       # EMBEDDING, id 1
    face = b"\x00" + bstr("함초롬돋움")

    def border_fill(bg_colorref: int | None) -> bytes:
        p = bytearray(44)
        if bg_colorref is not None:
            struct.pack_into("<I", p, 32, 0x1)            # fillflags: colorpattern
            struct.pack_into("<I", p, 36, bg_colorref)    # background COLORREF
        return bytes(p)

    def char_shape(size_pt: float, flags: int, color_bgr: int) -> bytes:
        p = bytearray(72)
        struct.pack_into("<H", p, 0, 0)  # ko face id 0
        struct.pack_into("<i", p, 42, int(size_pt * 100))
        struct.pack_into("<I", p, 46, flags)
        struct.pack_into("<I", p, 52, color_bgr)
        return bytes(p)

    def para_shape(align: int) -> bytes:
        return struct.pack("<I", align << 2) + b"\x00" * 40

    docinfo = b"".join([
        _rec(0x11, 0, id_mappings),
        _rec(0x12, 0, bin_data),
        _rec(0x13, 0, face),
        _rec(0x14, 0, border_fill(None)),          # borderfill id 1
        _rec(0x14, 0, border_fill(0x00CCFF)),      # id 2 — 배경 #FFCC00
        _rec(0x15, 0, char_shape(10.0, 0x0, 0x000000)),        # 0 보통
        _rec(0x15, 0, char_shape(12.0, 0x2 << 2, 0x0000FF)),   # 1 취소선+빨강
        _rec(0x15, 0, char_shape(14.0, 0x2, 0x000000)),        # 2 굵게
        _rec(0x19, 0, para_shape(1)),  # parashape 0: left
        _rec(0x19, 0, para_shape(3)),  # parashape 1: center
    ])

    # ── Section0 ──
    def para_header(parashape_id: int) -> bytes:
        p = bytearray(16)
        struct.pack_into("<H", p, 8, parashape_id)
        return bytes(p)

    def cell_props(col: int, row: int, colspan: int, rowspan: int,
                   w: int, h: int, borderfill: int) -> bytes:
        p = bytearray(40)
        struct.pack_into("<4H", p, 8, col, row, colspan, rowspan)
        struct.pack_into("<2i", p, 16, w, h)
        struct.pack_into("<H", p, 32, borderfill)
        return bytes(p)

    def cell(level: int, props: bytes, text: str) -> bytes:
        return b"".join([
            _rec(0x48, level, props),
            _rec(0x42, level + 1, para_header(0)),
            _rec(0x43, level + 2, _utf16(text)),
            _rec(0x44, level + 2, struct.pack("<II", 0, 0)),
        ])

    page = struct.pack("<6I", 59528, 84188, 8504, 8504, 5668, 4252) + b"\x00" * 16
    table_rec = struct.pack("<I", 0) + struct.pack("<HH", 2, 3) + b"\x00" * 16
    shape_comp = bytearray(36)
    struct.pack_into("<2i", shape_comp, 28, 14400, 7200)  # 2in × 1in
    shape_pic = bytearray(76)
    struct.pack_into("<H", shape_pic, 71, 1)  # bindata_id 1

    section = b"".join([
        _rec(0x49, 1, page),
        # 앵커 문단 — 가운데 정렬(parashape 1) + 취소선 런(charshape 1)
        _rec(0x42, 0, para_header(1)),
        _rec(0x43, 1, _utf16("가운데 취소선")),
        _rec(0x44, 1, struct.pack("<II", 0, 1)),
        # 표 2×3: (0,0) colspan2 배경, (0,2) rowspan2, (1,0), (1,1)
        _rec(0x47, 1, b" lbt" + b"\x00" * 8),
        _rec(0x4D, 2, table_rec),
        cell(2, cell_props(0, 0, 2, 1, 16000, 2000, 2), "병합 머리"),
        cell(2, cell_props(2, 0, 1, 2, 8000, 2000, 1), "세로 병합"),
        cell(2, cell_props(0, 1, 1, 1, 8000, 2000, 1), "좌"),
        cell(2, cell_props(1, 1, 1, 1, 8000, 2000, 1), "우"),
        # 그림 개체 (gso) — BinData/BIN0001.png
        _rec(0x47, 1, b" osg" + b"\x00" * 8),
        _rec(0x4C, 2, bytes(shape_comp)),
        _rec(0x55, 3, bytes(shape_pic)),
        # 머리말
        _rec(0x47, 1, b"daeh" + b"\x00" * 8),
        _rec(0x48, 2, b"\x00" * 8),
        _rec(0x42, 3, para_header(0)),
        _rec(0x43, 4, _utf16("머리말 텍스트")),
        # 본문 두 번째 문단 — 굵게(charshape 2), 글꼴 참조
        _rec(0x42, 0, para_header(0)),
        _rec(0x43, 1, _utf16("본문 끝")),
        _rec(0x44, 1, struct.pack("<II", 0, 2)),
    ])

    def comp(b: bytes) -> bytes:
        co = zlib.compressobj(6, zlib.DEFLATED, -15)
        return co.compress(b) + co.flush()

    return build_cfb({
        "FileHeader": bytes(header),
        "DocInfo": comp(docinfo),
        "BodyText/Section0": comp(section),
        "BinData/BIN0001.png": comp(_tiny_png()),
    })


def make_hwpx(text: str = "HWPX 본문", cell: str = "표셀") -> bytes:
    header_xml = """<?xml version="1.0" encoding="UTF-8"?>
<hh:head xmlns:hh="http://www.hancom.co.kr/hwpml/2011/head">
  <hh:refList>
    <hh:charProperties>
      <hh:charPr id="0" height="1000"/>
      <hh:charPr id="1" height="1600" textColor="#FF0000"><hh:bold/></hh:charPr>
    </hh:charProperties>
    <hh:paraProperties>
      <hh:paraPr id="0"><hh:align horizontal="CENTER"/></hh:paraPr>
    </hh:paraProperties>
  </hh:refList>
</hh:head>"""
    section_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section"
        xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p paraPrIDRef="0">
    <hp:run charPrIDRef="1">
      <hp:secPr><hp:pagePr width="59528" height="84188">
        <hp:margin left="8504" right="8504" top="5668" bottom="4252"/>
      </hp:pagePr></hp:secPr>
      <hp:t>{text}</hp:t>
    </hp:run>
  </hp:p>
  <hp:p>
    <hp:run charPrIDRef="0">
      <hp:tbl rowCnt="1" colCnt="1">
        <hp:tr><hp:tc>
          <hp:subList><hp:p><hp:run charPrIDRef="0"><hp:t>{cell}</hp:t></hp:run></hp:p></hp:subList>
          <hp:cellAddr colAddr="0" rowAddr="0"/><hp:cellSpan colSpan="1" rowSpan="1"/>
        </hp:tc></hp:tr>
      </hp:tbl>
    </hp:run>
  </hp:p>
</hs:sec>"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("mimetype", "application/hwp+zip")
        zf.writestr("Contents/header.xml", header_xml)
        zf.writestr("Contents/section0.xml", section_xml)
    return buf.getvalue()


def make_doc(text_paras: list[str]) -> bytes:
    body = "\r".join(text_paras) + "\r"
    text_bytes = _utf16(body)
    word = bytearray(0x0800)
    struct.pack_into("<H", word, 0, 0xA5EC)
    struct.pack_into("<H", word, 0x0A, 0x0000)  # 0Table
    text_fc = 0x0400
    word[text_fc:text_fc + len(text_bytes)] = text_bytes

    n = 1
    plc = struct.pack("<2I", 0, len(body))  # CP: 0..len
    plc += struct.pack("<HIH", 0, text_fc, 0)  # PCD: 유니코드(fc 그대로)
    pcdt = b"\x02" + struct.pack("<I", len(plc)) + plc
    table = bytearray(0x0400)
    fc_clx = 0x0040
    table[fc_clx:fc_clx + len(pcdt)] = pcdt
    struct.pack_into("<II", word, 0x01A2, fc_clx, len(pcdt))

    return build_cfb({"WordDocument": bytes(word), "0Table": bytes(table)})


def make_xls(strings: dict[tuple[int, int], str],
             numbers: dict[tuple[int, int], float]) -> bytes:
    def rec(rtype: int, payload: bytes) -> bytes:
        return struct.pack("<HH", rtype, len(payload)) + payload

    sst_list = sorted(set(strings.values()))
    sst_index = {s: i for i, s in enumerate(sst_list)}
    sst_payload = struct.pack("<II", len(strings), len(sst_list))
    for s in sst_list:
        raw = s.encode("utf-16le")
        sst_payload += struct.pack("<HB", len(s), 0x01) + raw

    font = struct.pack("<HHHH", 200, 0, 0x7FFF, 400) + b"\x00" * 6 + b"\x05\x00" + "Arial".encode("utf-16le")
    xf = struct.pack("<HH", 0, 0) + b"\x00" * 16

    globals_part = rec(0x0809, struct.pack("<HH", 0x0600, 0x0005) + b"\x00" * 12)
    globals_part += rec(0x0031, font)
    for _ in range(16):
        globals_part += rec(0x00E0, xf)
    globals_part += rec(0x00FC, sst_payload)

    sheet_cells = b""
    for (r, c), s in sorted(strings.items()):
        sheet_cells += rec(0x00FD, struct.pack("<HHHI", r, c, 15, sst_index[s]))
    for (r, c), v in sorted(numbers.items()):
        sheet_cells += rec(0x0203, struct.pack("<HHH", r, c, 15) + struct.pack("<d", v))
    sheet = rec(0x0809, struct.pack("<HH", 0x0600, 0x0010) + b"\x00" * 12)
    sheet += sheet_cells + rec(0x000A, b"")

    # BOUNDSHEET 의 BOF 오프셋은 전역부 길이에 의존 — 자리표시자로 한 번
    # 조립해 길이를 확정한 뒤 실제 오프셋으로 다시 쓴다 (2-패스).
    name = "Sheet1"

    def build_bs(bof: int) -> bytes:
        return rec(0x0085, struct.pack("<IBB", bof, 0, 0)
                   + bytes([len(name), 0]) + name.encode("latin-1"))

    eof = rec(0x000A, b"")
    bof_pos = len(globals_part) + len(build_bs(0)) + len(eof)
    workbook = globals_part + build_bs(bof_pos) + eof + sheet
    assert workbook[bof_pos:bof_pos + 2] == struct.pack("<H", 0x0809), "BOF offset mismatch"

    return build_cfb({"Workbook": workbook})


def make_ppt(slides: list[tuple[str, str]]) -> bytes:
    def rec(rtype: int, payload: bytes, ver: int = 0) -> bytes:
        return struct.pack("<HHI", ver, rtype, len(payload)) + payload

    doc_atom = rec(1001, struct.pack("<ii", 5760, 4320) + b"\x00" * 32)

    slwt_children = b""
    for title, body in slides:
        slwt_children += rec(1011, struct.pack("<II", 0, 0) + b"\x00" * 12)
        slwt_children += rec(3999, struct.pack("<I", 0))
        slwt_children += rec(4000, title.encode("utf-16le"))
        slwt_children += rec(3999, struct.pack("<I", 1))
        slwt_children += rec(4000, body.encode("utf-16le"))
    slwt = rec(4080, slwt_children, ver=0xF)
    document = rec(1000, doc_atom + slwt, ver=0xF)
    return build_cfb({"PowerPoint Document": document})


# ── 심화 픽스처: doc (CHPX/PAPX FKP + 표) ─────────────────────


def _chp_fkp(runs: list[tuple[int, bytes]], fc_end: int) -> bytes:
    """CHP FKP 페이지 — [(fc_start, grpprl)] + 마지막 경계."""
    crun = len(runs)
    page = bytearray(512)
    rgfc = [fc for fc, _ in runs] + [fc_end]
    struct.pack_into(f"<{crun + 1}I", page, 0, *rgfc)
    cursor = 510
    for i, (_fc, grpprl) in enumerate(runs):
        if not grpprl:
            continue
        size = 1 + len(grpprl)
        cursor -= size
        if cursor % 2:
            cursor -= 1
        page[cursor] = len(grpprl)
        page[cursor + 1:cursor + 1 + len(grpprl)] = grpprl
        page[(crun + 1) * 4 + i] = cursor // 2
    page[511] = crun
    return bytes(page)


def _pap_fkp(runs: list[tuple[int, bytes]], fc_end: int) -> bytes:
    """PAP FKP 페이지 — BX 13B, papx = cb(1B[+1B]) + istd(2B) + sprms."""
    crun = len(runs)
    page = bytearray(512)
    rgfc = [fc for fc, _ in runs] + [fc_end]
    struct.pack_into(f"<{crun + 1}I", page, 0, *rgfc)
    cursor = 510
    for i, (_fc, sprms) in enumerate(runs):
        if not sprms:
            continue
        body = b"\x00\x00" + sprms  # istd 0 + sprms
        t = len(body)
        if t % 2:  # size = 2*cb - 1
            blob = bytes([(t + 1) // 2]) + body
        else:      # size = 2*cb2 (첫 바이트 0)
            blob = bytes([0, t // 2]) + body
        cursor -= len(blob)
        if cursor % 2:
            cursor -= 1
        page[cursor:cursor + len(blob)] = blob
        page[(crun + 1) * 4 + i * 13] = cursor // 2
    page[511] = crun
    return bytes(page)


def make_doc_rich() -> bytes:
    """정렬/런 스타일/글꼴/표(fInTable·fTtp)를 전부 싣는 Word97 픽스처."""
    body = (
        "제목 가운데\r"          # cp 0-5, mark @6 — 가운데 정렬
        "굵은빨강취소\r"          # cp 7-12, mark @13 — 앞 3자 bold+red, 뒤 3자 strike
        "A1\x07B1\x07\x07"       # 표 1행: 셀 2 + 행마크
        "A2\x07B2\x07\x07"       # 표 2행
        "끝문단\r"
    )
    text_bytes = _utf16(body)
    word = bytearray(0x1000)
    struct.pack_into("<H", word, 0, 0xA5EC)
    struct.pack_into("<H", word, 0x0A, 0x0000)  # 0Table
    text_fc = 0x0400
    word[text_fc:text_fc + len(text_bytes)] = text_bytes

    def fc(cp: int) -> int:
        return text_fc + 2 * cp

    # CHPX: [0,7) 기본 / [7,10) bold+red / [10,13) strike+underline+14pt+글꼴0
    sprm_bold = struct.pack("<HB", 0x0835, 1)
    sprm_red = struct.pack("<H", 0x6870) + bytes([255, 0, 0, 0])
    sprm_strike = struct.pack("<HB", 0x0837, 1)
    sprm_kul = struct.pack("<HB", 0x2A3E, 1)
    sprm_hps = struct.pack("<HH", 0x4A43, 28)  # 14pt
    sprm_ftc = struct.pack("<HH", 0x4A4F, 0)
    chp_page = _chp_fkp([
        (fc(0), b""),
        (fc(7), sprm_bold + sprm_red),
        (fc(10), sprm_strike + sprm_kul + sprm_hps + sprm_ftc),
        (fc(13), b""),
    ], fc(len(body)))

    # PAPX: 문단 마크 fc 로 구간을 나눈다. 표 구간 문단은 fInTable,
    # 행마크(\x07 단독)는 fInTable+fTtp. 첫 문단은 jc=center.
    sprm_jc = struct.pack("<HB", 0x2403, 1)
    sprm_intbl = struct.pack("<HB", 0x2416, 1)
    sprm_ttp = struct.pack("<HB", 0x2417, 1)
    # cp 배치: "제목 가운데\r"=0..6, "굵은빨강취소\r"=7..13,
    # A1@14-15 \x07@16, B1@17-18 \x07@19, 행마크@20,
    # A2@21-22 \x07@23, B2@24-25 \x07@26, 행마크@27, "끝문단\r"=28..31
    pap_page = _pap_fkp([
        (fc(0), sprm_jc),            # 문단1 (mark @6)
        (fc(7), b""),                # 문단2
        (fc(14), sprm_intbl),        # A1 셀
        (fc(17), sprm_intbl),        # B1 셀
        (fc(20), sprm_intbl + sprm_ttp),  # 행마크1
        (fc(21), sprm_intbl),        # A2
        (fc(24), sprm_intbl),        # B2
        (fc(27), sprm_intbl + sprm_ttp),  # 행마크2
        (fc(28), b""),               # 끝문단
    ], fc(len(body)))

    # FKP 페이지는 512 정렬 오프셋에 놓인다 — pn 6/7 사용 (word 0x1000 안)
    word[6 * 512:7 * 512] = chp_page
    word[7 * 512:8 * 512] = pap_page

    table = bytearray(0x0800)
    # Clx/PlcPcd
    plc = struct.pack("<2I", 0, len(body)) + struct.pack("<HIH", 0, text_fc, 0)
    pcdt = b"\x02" + struct.pack("<I", len(plc)) + plc
    fc_clx = 0x0040
    table[fc_clx:fc_clx + len(pcdt)] = pcdt
    struct.pack_into("<II", word, 0x01A2, fc_clx, len(pcdt))
    # PlcfBteChpx / PlcfBtePapx — 구간 1개, FKP pn 6/7
    bte_chp = struct.pack("<2I", fc(0), fc(len(body))) + struct.pack("<I", 6)
    bte_pap = struct.pack("<2I", fc(0), fc(len(body))) + struct.pack("<I", 7)
    fc_bte_chp, fc_bte_pap = 0x0200, 0x0240
    table[fc_bte_chp:fc_bte_chp + len(bte_chp)] = bte_chp
    table[fc_bte_pap:fc_bte_pap + len(bte_pap)] = bte_pap
    struct.pack_into("<II", word, 0x00FA, fc_bte_chp, len(bte_chp))
    struct.pack_into("<II", word, 0x0102, fc_bte_pap, len(bte_pap))
    # SttbfFfn — 글꼴 1개 "바탕"
    name = "바탕".encode("utf-16le") + b"\x00\x00"
    ffn_body = bytes(39) + name
    sttb = struct.pack("<HHH", 0xFFFF, 1, 0) + bytes([len(ffn_body)]) + ffn_body
    fc_ffn = 0x0300
    table[fc_ffn:fc_ffn + len(sttb)] = sttb
    struct.pack_into("<II", word, 0x0112, fc_ffn, len(sttb))

    return build_cfb({"WordDocument": bytes(word), "0Table": bytes(table)})


# ── 심화 픽스처: xls (스타일/정렬/채우기/팔레트) ──────────────


def make_xls_rich() -> bytes:
    def rec(rtype: int, payload: bytes) -> bytes:
        return struct.pack("<HH", rtype, len(payload)) + payload

    s_text = "스타일셀"
    sst_payload = struct.pack("<II", 1, 1)
    sst_payload += struct.pack("<HB", len(s_text), 0x01) + s_text.encode("utf-16le")

    def font(height: int, grbit: int, icv: int, weight: int, uls: int) -> bytes:
        return (struct.pack("<HHHH", height, grbit, icv, weight)
                + struct.pack("<HBBBB", 0, uls, 0, 0, 0)
                + b"\x05\x01" + "Arial".encode("utf-16le"))

    # 폰트 0-3 기본, (인덱스 4 없음 규약) → 5번째 폰트가 ifnt=5
    fonts = [font(200, 0, 0x7FFF, 400, 0)] * 4
    fonts.append(font(240, 0x0008, 40, 700, 1))  # 12pt bold strike underline icv40

    def xf(ifnt: int, alc: int, pattern_icv: int | None) -> bytes:
        p = bytearray(20)
        struct.pack_into("<HH", p, 0, ifnt, 0)
        p[6] = alc
        if pattern_icv is not None:
            struct.pack_into("<i", p, 14, 1 << 26)      # solid
            struct.pack_into("<H", p, 18, pattern_icv)  # 전경색 icv
        return bytes(p)

    globals_part = rec(0x0809, struct.pack("<HH", 0x0600, 0x0005) + b"\x00" * 12)
    for f in fonts:
        globals_part += rec(0x0031, f)
    # PALETTE — icv 40 을 (255, 204, 0) 으로 재정의
    pal = struct.pack("<H", 56)
    for k in range(56):
        pal += bytes([255, 204, 0, 0]) if k == 32 else bytes([k, k, k, 0])
    globals_part += rec(0x0092, pal)
    for _ in range(15):
        globals_part += rec(0x00E0, xf(0, 0, None))
    globals_part += rec(0x00E0, xf(5, 0x02, 40))  # ixfe 15: 가운데+채움+스타일폰트
    globals_part += rec(0x00FC, sst_payload)

    sheet_cells = rec(0x00FD, struct.pack("<HHHI", 0, 0, 15, 0))
    sheet = rec(0x0809, struct.pack("<HH", 0x0600, 0x0010) + b"\x00" * 12)
    sheet += sheet_cells + rec(0x000A, b"")

    name = "S"

    def build_bs(bof: int) -> bytes:
        return rec(0x0085, struct.pack("<IBB", bof, 0, 0)
                   + bytes([len(name), 0]) + name.encode("latin-1"))

    eof = rec(0x000A, b"")
    bof_pos = len(globals_part) + len(build_bs(0)) + len(eof)
    workbook = globals_part + build_bs(bof_pos) + eof + sheet
    return build_cfb({"Workbook": workbook})


# ── 심화 픽스처: ppt (StyleTextPropAtom 런 스타일) ────────────


def make_ppt_rich() -> bytes:
    def rec(rtype: int, payload: bytes, ver: int = 0) -> bytes:
        return struct.pack("<HHI", ver, rtype, len(payload)) + payload

    doc_atom = rec(1001, struct.pack("<ii", 5760, 4320) + b"\x00" * 32)

    body_text = "굵은줄\r빨강취소"
    # 문단 런: 전체 + 1 자 — alignment(0x800)=center(1)
    para_run = (struct.pack("<I", len(body_text) + 1) + struct.pack("<H", 0)
                + struct.pack("<I", 0x800) + struct.pack("<H", 1))
    # 문자 런 1: "굵은줄\r"(4자) — bold + size 24
    char_run1 = (struct.pack("<I", 4)
                 + struct.pack("<I", 0x0001 | 0x20000)
                 + struct.pack("<H", 0x0001) + struct.pack("<H", 24))
    # 문자 런 2: "빨강취소"(4자)+1 — strike + color 빨강
    char_run2 = (struct.pack("<I", 5)
                 + struct.pack("<I", 0x0100 | 0x40000)
                 + struct.pack("<H", 0x0100) + bytes([255, 0, 0, 0]))
    style = rec(4001, para_run + char_run1 + char_run2)

    slwt_children = rec(1011, struct.pack("<II", 0, 0) + b"\x00" * 12)
    slwt_children += rec(3999, struct.pack("<I", 0))
    slwt_children += rec(4000, "제목".encode("utf-16le"))
    slwt_children += rec(3999, struct.pack("<I", 1))
    slwt_children += rec(4000, body_text.encode("utf-16le")) + style
    slwt = rec(4080, slwt_children, ver=0xF)
    document = rec(1000, doc_atom + slwt, ver=0xF)
    return build_cfb({"PowerPoint Document": document})


# ─────────────────────────────────────────────────────────────
# 테스트
# ─────────────────────────────────────────────────────────────


def _render_svg_pages(tmp_path: Path, name: str, content: bytes) -> str:
    src = tmp_path / name
    src.write_bytes(content)
    out = tmp_path / "render"
    result = render_doc(src, to="svg", out_dir=out)
    assert result.page_count >= 1
    return "".join(p.read_text(encoding="utf-8") for p in result.paths)


class TestHwp:
    def test_hwp_renders_text_styles_and_table(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "sample.hwp", make_hwp())
        assert "한글 렌더 검증" in svg
        assert "굵은 부분" in svg
        assert "셀A" in svg and "셀B" in svg
        # charshape 1 = 굵게 + 빨강 — 스타일이 살아서 렌더까지 도달한다
        assert 'font-weight="bold"' in svg
        assert "#FF0000" in svg or "FF0000" in svg

    def test_password_protected_is_refused_honestly(self, tmp_path):
        content = bytearray(make_hwp())
        # FileHeader 는 첫 스트림 — flags 에 password 비트를 세운다
        idx = content.find(b"HWP Document File")
        struct.pack_into("<I", content, idx + 36, 0x3)
        with pytest.raises(LegacyConvertError):
            convert_to_ooxml(bytes(content), "hwp")

    def test_non_hwp_bytes_are_refused(self):
        with pytest.raises(LegacyConvertError):
            convert_to_ooxml(b"not an ole file at all", "hwp")


@pytest.fixture(scope="module")
def rich_docx():
    from docx import Document

    docx_bytes, fmt = convert_to_ooxml(make_hwp_rich(), "hwp")
    assert fmt == "docx"
    return Document(io.BytesIO(docx_bytes))


class TestHwpFidelity:
    """심화 충실도 — 표 75 병합/배경, 표 38 정렬, 표 28 취소선·글꼴,
    표 102 그림, 머리말. 픽스처는 make_hwp_rich (실파일 레벨 배치)."""

    def test_colspan_merge(self, rich_docx):
        tbl = rich_docx.tables[0]
        assert tbl.cell(0, 0).text.strip() == "병합 머리"
        xml = tbl._tbl.xml
        assert 'gridSpan' in xml and 'w:val="2"' in xml
        # 병합 셀과 (0,1) 이 같은 tc 를 공유한다
        assert tbl.cell(0, 0)._tc is tbl.cell(0, 1)._tc

    def test_rowspan_merge(self, rich_docx):
        tbl = rich_docx.tables[0]
        assert tbl.cell(0, 2).text.strip() == "세로 병합"
        assert 'vMerge' in tbl._tbl.xml
        assert tbl.cell(0, 2)._tc is tbl.cell(1, 2)._tc

    def test_cell_background_from_borderfill(self, rich_docx):
        # borderfill id 2 = COLORREF 0x00CCFF → #FFCC00
        assert 'FFCC00' in rich_docx.tables[0]._tbl.xml

    def test_paragraph_alignment(self, rich_docx):
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        first = rich_docx.paragraphs[0]
        assert first.text == "가운데 취소선"
        assert first.alignment == WD_ALIGN_PARAGRAPH.CENTER

    def test_strike_color_and_face(self, rich_docx):
        run = rich_docx.paragraphs[0].runs[0]
        assert run.font.strike
        assert str(run.font.color.rgb) == "FF0000"
        assert run.font.name == "함초롬돋움"
        last = [p for p in rich_docx.paragraphs if p.text == "본문 끝"][0]
        assert last.runs[0].bold

    def test_image_embedded_with_size(self, rich_docx):
        shapes = rich_docx.inline_shapes
        assert len(shapes) == 1
        # 14400×7200 HWPUNIT = 2in × 1in
        assert abs(shapes[0].width - 914400 * 2) < 2000
        assert abs(shapes[0].height - 914400) < 2000

    def test_header_text(self, rich_docx):
        assert "머리말 텍스트" in rich_docx.sections[0].header.paragraphs[0].text

    def test_column_widths_reach_grid(self, rich_docx):
        xml = rich_docx.tables[0]._tbl.xml
        assert "gridCol" in xml

    def test_e2e_svg_render(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "rich.hwp", make_hwp_rich())
        assert "병합 머리" in svg and svg.count("병합 머리") == 1
        assert "세로 병합" in svg
        assert "#FFCC00" in svg          # 셀 배경
        assert "line-through" in svg     # 취소선
        assert "<image" in svg           # 그림 개체
        assert "머리말 텍스트" in svg    # 머리말


class TestHwpx:
    def test_hwpx_renders_styles_table_and_page_geometry(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "sample.hwpx", make_hwpx())
        assert "HWPX 본문" in svg
        assert "표셀" in svg
        assert 'font-weight="bold"' in svg   # charPr 1 bold
        assert "FF0000" in svg               # textColor

    def test_zip_without_sections_is_refused(self):
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("mimetype", "application/hwp+zip")
        with pytest.raises(LegacyConvertError):
            convert_to_ooxml(buf.getvalue(), "hwpx")


class TestDoc:
    def test_doc_pieces_reassemble_paragraphs(self, tmp_path):
        svg = _render_svg_pages(
            tmp_path, "sample.doc",
            make_doc(["Word 문서 첫 문단", "둘째 문단 내용"]),
        )
        assert "Word 문서 첫 문단" in svg
        assert "둘째 문단 내용" in svg

    def test_doc_field_instruction_is_dropped_result_kept(self, tmp_path):
        # \x13 명령 \x14 결과 \x15 — 결과만 남아야 한다
        raw = "앞 \x13PAGEREF _Toc\x14보이는 결과\x15 뒤"
        svg = _render_svg_pages(tmp_path, "f.doc", make_doc([raw]))
        assert "보이는 결과" in svg
        assert "PAGEREF" not in svg


@pytest.fixture(scope="module")
def rich():
    from docx import Document

    docx_bytes, fmt = convert_to_ooxml(make_doc_rich(), "doc")
    assert fmt == "docx"
    return Document(io.BytesIO(docx_bytes))


class TestDocFidelity:
    """CHPX/PAPX FKP — 런 스타일/정렬/글꼴/표(fInTable·fTtp)."""

    def test_alignment_center(self, rich):
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        first = rich.paragraphs[0]
        assert first.text == "제목 가운데"
        assert first.alignment == WD_ALIGN_PARAGRAPH.CENTER

    def test_char_runs_split_by_chpx(self, rich):
        para = rich.paragraphs[1]
        assert para.text == "굵은빨강취소"
        runs = [r for r in para.runs if r.text]
        assert runs[0].text == "굵은빨강"[:3]
        assert runs[0].bold
        assert str(runs[0].font.color.rgb) == "FF0000"
        last = runs[-1]
        assert last.font.strike and last.underline
        assert last.font.size.pt == 14.0
        assert last.font.name == "바탕"

    def test_table_from_cell_and_row_marks(self, rich):
        assert len(rich.tables) == 1
        tbl = rich.tables[0]
        assert len(tbl.rows) == 2 and len(tbl.columns) == 2
        assert tbl.cell(0, 0).text.strip() == "A1"
        assert tbl.cell(0, 1).text.strip() == "B1"
        assert tbl.cell(1, 0).text.strip() == "A2"
        assert tbl.cell(1, 1).text.strip() == "B2"
        # 표 뒤 본문이 표에 빨려들지 않는다
        assert any(p.text == "끝문단" for p in rich.paragraphs)

    def test_e2e_svg(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "rich.doc", make_doc_rich())
        assert "제목 가운데" in svg
        assert 'font-weight="bold"' in svg
        assert "#FF0000" in svg
        assert "line-through" in svg
        assert "A1" in svg and "B2" in svg  # 표가 격자로 렌더


@pytest.fixture(scope="module")
def rich_ws():
    from openpyxl import load_workbook

    xlsx_bytes, fmt = convert_to_ooxml(make_xls_rich(), "xls")
    assert fmt == "xlsx"
    return load_workbook(io.BytesIO(xlsx_bytes)).active


class TestXlsFidelity:
    """FONT 색/밑줄/취소선 + XF 정렬/solid 채우기 + PALETTE 재정의."""


    def test_font_styles_and_palette_color(self, rich_ws):
        cell = rich_ws["A1"]
        assert cell.value == "스타일셀"
        assert cell.font.bold and cell.font.strike
        assert cell.font.underline == "single"
        assert cell.font.size == 12.0
        # PALETTE 재정의: icv 40 → FFCC00
        assert (cell.font.color.rgb or "").endswith("FFCC00")

    def test_alignment_and_fill(self, rich_ws):
        cell = rich_ws["A1"]
        assert cell.alignment.horizontal == "center"
        assert cell.fill.patternType == "solid"
        assert (cell.fill.fgColor.rgb or "").endswith("FFCC00")

    def test_e2e_svg(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "rich.xls", make_xls_rich())
        assert "스타일셀" in svg
        assert "#FFCC00" in svg  # 채우기가 렌더에 도달


@pytest.fixture(scope="module")
def rich_slide():
    from pptx import Presentation

    pptx_bytes, fmt = convert_to_ooxml(make_ppt_rich(), "ppt")
    assert fmt == "pptx"
    return Presentation(io.BytesIO(pptx_bytes)).slides[0]


class TestPptFidelity:
    """StyleTextPropAtom — 문자 런 스타일/문단 정렬."""


    def _body_paras(self, slide):
        boxes = [sh for sh in slide.shapes if sh.has_text_frame]
        body = boxes[-1]
        return body.text_frame.paragraphs

    def test_char_run_styles(self, rich_slide):
        paras = self._body_paras(rich_slide)
        r1 = paras[0].runs[0]
        assert r1.text == "굵은줄"
        assert r1.font.bold
        assert r1.font.size.pt == 24.0
        r2 = paras[1].runs[0]
        assert r2.text == "빨강취소"
        assert str(r2.font.color.rgb) == "FF0000"
        assert r2.font._rPr.get("strike") == "sngStrike"

    def test_paragraph_alignment(self, rich_slide):
        from pptx.enum.text import PP_ALIGN

        paras = self._body_paras(rich_slide)
        assert paras[0].alignment == PP_ALIGN.CENTER

    def test_e2e_svg(self, tmp_path):
        svg = _render_svg_pages(tmp_path, "rich.ppt", make_ppt_rich())
        assert "굵은줄" in svg and "빨강취소" in svg


class TestXls:
    def test_xls_strings_numbers_render_in_grid(self, tmp_path):
        content = make_xls(
            strings={(0, 0): "항목", (0, 1): "금액"},
            numbers={(1, 1): 1234.5},
        )
        svg = _render_svg_pages(tmp_path, "sample.xls", content)
        assert "항목" in svg and "금액" in svg
        assert "1234.5" in svg

    def test_biff5_is_refused(self):
        biff5 = struct.pack("<HH", 0x0809, 8) + struct.pack("<HH", 0x0500, 0x0005) + b"\x00" * 4
        with pytest.raises(LegacyConvertError):
            convert_to_ooxml(build_cfb({"Workbook": biff5}), "xls")


class TestPpt:
    def test_ppt_slides_titles_and_bodies(self, tmp_path):
        content = make_ppt([
            ("첫 슬라이드 제목", "본문 한 줄"),
            ("둘째 제목", "둘째 본문"),
        ])
        svg = _render_svg_pages(tmp_path, "sample.ppt", content)
        assert "첫 슬라이드 제목" in svg
        assert "둘째 본문" in svg

    def test_slide_count_preserved(self, tmp_path):
        content = make_ppt([("A", "a"), ("B", "b"), ("C", "c")])
        src = tmp_path / "three.ppt"
        src.write_bytes(content)
        result = render_doc(src, to="svg", out_dir=tmp_path / "r")
        assert result.page_count == 3
