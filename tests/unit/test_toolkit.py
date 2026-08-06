"""Tests for the agent-toolkit surfaces: facade, agent_tools, local MCP.

LLM-backed verbs (generate/edit) are exercised with a stubbed client via
the existing edit_deck stubs elsewhere; here we cover the deterministic
verbs end-to-end plus schema/dispatch/lazy-import contracts.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

import xgen_edit2docs
from xgen_edit2docs import analyze_pptx, preview_pptx, run_tool, set_pptx_text
from xgen_edit2docs.agent_tools import ANTHROPIC_TOOLS, TOOL_NAMES


@pytest.fixture
def deck_path(tmp_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "원래 제목"
    gf = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(2))
    gf.table.cell(0, 0).text = "셀 텍스트"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


class TestLazyPackageSurface:
    def test_version_and_lazy_exports(self):
        assert xgen_edit2docs.__version__
        assert callable(xgen_edit2docs.generate_pptx)
        assert callable(xgen_edit2docs.edit_pptx)
        assert "ANTHROPIC_TOOLS" in dir(xgen_edit2docs)

    def test_unknown_attribute_raises(self):
        with pytest.raises(AttributeError):
            xgen_edit2docs.does_not_exist  # noqa: B018


class TestFacadeDeterministicVerbs:
    def test_preview_strings_and_files(self, deck_path, tmp_path):
        svgs = preview_pptx(deck_path)
        assert len(svgs) == 1 and "원래 제목" in svgs[0]
        paths = preview_pptx(deck_path, out_dir=tmp_path / "svg")
        assert [p.name for p in paths] == ["slide_000.svg"]
        assert "원래 제목" in Path(paths[0]).read_text(encoding="utf-8")

    def test_analyze_lists_addressable_paragraphs(self, deck_path):
        info = analyze_pptx(deck_path)
        assert info["page_count"] == 1
        texts = info["slides"][0]["texts"]
        shape_entries = [t for t in texts if "shape_id" in t]
        table_entries = [t for t in texts if "table_id" in t]
        assert any(t["text"] == "원래 제목" for t in shape_entries)
        cell = next(t for t in table_entries if t["text"] == "셀 텍스트")
        assert (cell["row"], cell["col"]) == (0, 0)

    def test_set_text_via_analyze_addresses(self, deck_path, tmp_path):
        info = analyze_pptx(deck_path)
        target = next(
            t for t in info["slides"][0]["texts"] if t.get("text") == "원래 제목"
        )
        out = tmp_path / "out.pptx"
        result = set_pptx_text(
            deck_path,
            [
                {
                    "slide": 0,
                    "shape_id": target["shape_id"],
                    "para": target["para"],
                    "new_text": "파사드 수정",
                    "old_text": target["text"],
                }
            ],
            output=out,
        )
        assert result.applied == 1 and result.path == out
        assert "파사드 수정" in preview_pptx(out)[0]

    def test_zero_applied_leaves_input_untouched(self, deck_path, tmp_path):
        result = set_pptx_text(
            deck_path,
            [{"slide": 0, "shape_id": 99999, "para": 0, "new_text": "x"}],
            output=tmp_path / "never.pptx",
        )
        assert result.applied == 0
        assert not (tmp_path / "never.pptx").exists()
        assert result.path == deck_path


class TestAgentTools:
    def test_schemas_are_anthropic_shaped(self):
        assert TOOL_NAMES == [
            "doc_guide", "analyze_doc", "render_doc", "set_doc_text",
            "arrange_doc", "read_doc_xml", "set_doc_xml", "build_doc",
            "generate_doc", "edit_doc",
        ]
        for tool in ANTHROPIC_TOOLS:
            assert set(tool) == {"name", "description", "input_schema"}
            schema = tool["input_schema"]
            assert schema["type"] == "object"
            if tool["name"] != "doc_guide":  # guide's only param is optional
                assert schema["required"], tool["name"]

    def test_descriptions_stay_compact(self):
        """Progressive disclosure contract: the frontmatter tier must stay
        small — fat how-to lives behind doc_guide(topic)."""
        for tool in ANTHROPIC_TOOLS:
            assert len(tool["description"]) <= 320, (
                f"{tool['name']} description grew to "
                f"{len(tool['description'])} chars — move detail into "
                "agent_guide GUIDES instead"
            )

    def test_openai_surface_matches(self):
        from xgen_edit2docs.agent_tools import OPENAI_TOOLS, tool_specs

        assert [t["function"]["name"] for t in OPENAI_TOOLS] == TOOL_NAMES
        for a, o in zip(ANTHROPIC_TOOLS, OPENAI_TOOLS):
            assert o["type"] == "function"
            assert o["function"]["description"] == a["description"]
            assert o["function"]["parameters"] == a["input_schema"]
        assert tool_specs("openai") is OPENAI_TOOLS
        assert tool_specs("anthropic") is ANTHROPIC_TOOLS
        with pytest.raises(ValueError):
            tool_specs("gemini")

    def test_dispatch_deterministic_tools(self, deck_path, tmp_path):
        info = run_tool("analyze_doc", {"doc": str(deck_path)})
        assert info["page_count"] == 1 and info["format"] == "pptx"
        res = run_tool(
            "render_doc",
            {"doc": str(deck_path), "to": "md", "out_dir": str(tmp_path / "s")},
        )
        assert res["page_count"] == 1 and Path(res["paths"][0]).exists()

    def test_unknown_tool_raises(self):
        with pytest.raises(ValueError, match="unknown xgen_edit2docs tool"):
            run_tool("rm_rf_slash", {})


class TestDocGuide:
    """The hierarchical skill guide — progressive disclosure entry point."""

    def test_root_map_splits_generate_and_edit(self):
        res = run_tool("doc_guide", {})
        assert res["topic"] == ""
        assert "GENERATE" in res["guide"] and "EDIT" in res["guide"]
        # every tool is discoverable from the root map
        for name in TOOL_NAMES:
            assert name in res["guide"], f"{name} missing from root map"
        assert res["topics"]

    def test_every_topic_resolves(self):
        from xgen_edit2docs.agent_guide import TOPICS

        for topic in TOPICS:
            res = run_tool("doc_guide", {"topic": topic})
            assert res["topic"] == topic
            assert len(res["guide"]) > 100, topic

    def test_parent_prefix_joins_children(self):
        res = run_tool("doc_guide", {"topic": "recipes"})
        assert "COPY / MOVE / DELETE a slide" in res["guide"]  # recipes.slides
        assert "RECOLOR" in res["guide"]  # recipes.colors

    def test_exact_topic_lists_subtopics(self):
        res = run_tool("doc_guide", {"topic": "edit"})
        assert "Subtopics:" in res["guide"]
        assert "edit.xml" in res["guide"]

    def test_unknown_topic_is_never_a_dead_end(self):
        res = run_tool("doc_guide", {"topic": "no-such-thing"})
        assert "GENERATE" in res["guide"]  # falls back to the map
        assert res["topics"]

    def test_host_name_mapping(self):
        """Hosts that rename tools (geny-executor) get the guide rendered
        with THEIR names."""
        from xgen_edit2docs.agent_guide import doc_guide

        names = {"analyze_doc": "DocAnalyze", "set_doc_xml": "DocXmlEdit",
                 "doc_guide": "DocGuide"}
        res = doc_guide("recipes.slides", names=names)
        assert "DocXmlEdit" in res["guide"]
        assert "set_doc_xml" not in res["guide"]


class TestBuildDoc:
    """Deterministic generation — generate_doc's engine without the model."""

    def test_build_docx_from_markdown(self, tmp_path):
        out = tmp_path / "r.docx"
        res = run_tool(
            "build_doc",
            {"spec": "# Report\n\nBody **text**.\n\n- a\n- b", "output": str(out)},
        )
        assert Path(res["path"]) == out and out.exists()
        info = xgen_edit2docs.analyze_doc(str(out))
        assert info["format"] == "docx"

    def test_build_xlsx_from_spec(self, tmp_path):
        out = tmp_path / "r.xlsx"
        res = run_tool(
            "build_doc",
            {
                "spec": {"sheets": [{"name": "S", "headers": ["x", "y"],
                                     "rows": [[1, 2], [3, 4]]}]},
                "output": str(out),
            },
        )
        assert res["page_count"] == 1 and out.exists()
        assert "3" in str(xgen_edit2docs.analyze_doc(str(out)))

    def test_build_pptx_from_slide_spec(self, tmp_path):
        out = tmp_path / "r.pptx"
        res = run_tool(
            "build_doc",
            {
                "spec": {"slides": [
                    {"layout": "title", "title": "Deck", "subtitle": "2026"},
                    {"layout": "content", "title": "Agenda",
                     "bullets": ["A", {"text": "A.1", "level": 1}], "notes": "n"},
                ]},
                "output": str(out),
            },
        )
        assert res["page_count"] == 2 and out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        assert prs.slides[0].shapes.title.text == "Deck"
        assert prs.slides[1].shapes.title.text == "Agenda"

    def test_build_pptx_rejects_wrong_spec_type(self, tmp_path):
        with pytest.raises(ValueError, match="slides|dict"):
            run_tool("build_doc", {"spec": "markdown", "output": str(tmp_path / "x.pptx")})

    def test_build_pptx_rejects_empty_slides(self, tmp_path):
        with pytest.raises(ValueError, match="slides"):
            run_tool("build_doc", {"spec": {"slides": []}, "output": str(tmp_path / "x.pptx")})


class TestDocXml:
    """Direct OOXML XML editing — the universal deterministic escape hatch."""

    @pytest.fixture
    def chart_deck(self, tmp_path: Path) -> Path:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C"]
        cd.add_series("S1", (1, 2, 3))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(4), cd,
        )
        p = tmp_path / "chart.pptx"
        prs.save(str(p))
        return p

    def test_list_parts_maps_the_package(self, chart_deck):
        res = run_tool("read_doc_xml", {"doc": str(chart_deck)})
        names = [p["part"] for p in res["parts"]]
        assert any("slides/slide1.xml" in n for n in names)
        assert any("charts/chart1.xml" in n for n in names)

    def test_read_one_part_returns_exact_xml(self, chart_deck):
        res = run_tool(
            "read_doc_xml", {"doc": str(chart_deck), "part": "ppt/charts/chart1.xml"}
        )
        assert "<c:ser>" in res["xml"] and "S1" in res["xml"]

    def test_recolor_chart_series_via_xml_patch(self, chart_deck, tmp_path):
        """THE user scenario: change bar color — impossible via edit_chart,
        trivial via a direct XML patch on the series properties."""
        out = tmp_path / "red.pptx"
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{
                    "find": "</c:tx>",
                    "replace": (
                        "</c:tx><c:spPr><a:solidFill>"
                        '<a:srgbClr val="FF0000"/>'
                        "</a:solidFill></c:spPr>"
                    ),
                }],
                "output": str(out),
            },
        )
        assert res["applied"] == 1, res
        # python-pptx sees the explicit red fill on the series.
        prs = Presentation(str(out))
        chart = next(
            s for sl in prs.slides for s in sl.shapes if s.has_chart
        ).chart
        fill = chart.series[0].format.fill
        assert str(fill.fore_color.rgb) == "FF0000"

    def test_malformed_result_is_rejected_not_written(self, chart_deck):
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{"find": "</c:chartSpace>", "replace": "<broken"}],
            },
        )
        assert res["applied"] == 0
        assert any(r["status"] == "invalid" for r in res["results"])
        # Nothing was written: the source still parses.
        assert Presentation(str(chart_deck))

    def test_not_found_reports_status(self, chart_deck):
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{"find": "<no-such-element/>", "replace": "x"}],
            },
        )
        assert res["applied"] == 0
        assert res["results"][0]["status"] == "not_found"

    def test_missing_part_suggests_names(self, chart_deck):
        with pytest.raises(ValueError, match="chart1.xml"):
            run_tool(
                "read_doc_xml", {"doc": str(chart_deck), "part": "chart1.xml"}
            )

    def test_add_slide_via_pure_xml_tool_calls(self, tmp_path):
        """The power proof: ADD A SLIDE with nothing but read/set_doc_xml —
        create slideN.xml + its rels, register the content type, patch
        presentation.xml + its rels. python-pptx must see 2 slides."""
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "First"
        doc = tmp_path / "deck.pptx"
        prs.save(str(doc))
        d = {"doc": str(doc)}

        # 1) copy slide1's xml + rels as the template for slide2
        slide_xml = run_tool(
            "read_doc_xml", {**d, "part": "ppt/slides/slide1.xml"}
        )["xml"].replace("First", "Second")
        rels_xml = run_tool(
            "read_doc_xml", {**d, "part": "ppt/slides/_rels/slide1.xml.rels"}
        )["xml"]
        # 2) create the new parts (content type registered for the slide)
        r = run_tool("set_doc_xml", {
            **d, "part": "ppt/slides/slide2.xml", "xml": slide_xml,
            "content_type": (
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.slide+xml"
            ),
            "output": str(doc),
        })
        assert r["applied"] == 1, r
        r = run_tool("set_doc_xml", {
            **d, "part": "ppt/slides/_rels/slide2.xml.rels", "xml": rels_xml,
            "output": str(doc),
        })
        assert r["applied"] == 1, r
        # 3) wire it into the presentation: rels + sldIdLst
        new_rel = (
            '<Relationship Id="rIdNew" Type="http://schemas.openxmlformats.org/'
            'officeDocument/2006/relationships/slide" Target="slides/slide2.xml"/>'
            "</Relationships>"
        )
        r = run_tool("set_doc_xml", {
            **d, "part": "ppt/_rels/presentation.xml.rels",
            "edits": [{"find": "</Relationships>", "replace": new_rel}],
            "output": str(doc),
        })
        assert r["applied"] == 1, r
        r = run_tool("set_doc_xml", {
            **d, "part": "ppt/presentation.xml",
            "edits": [{
                "find": "</p:sldIdLst>",
                "replace": '<p:sldId id="9999" r:id="rIdNew"/></p:sldIdLst>',
            }],
            "output": str(doc),
        })
        assert r["applied"] == 1, r
        # 4) acceptance: python-pptx opens it and sees both slides
        check = Presentation(str(doc))
        assert len(check.slides) == 2
        titles = [sl.shapes.title.text for sl in check.slides]
        assert titles == ["First", "Second"]

    def test_delete_part_roundtrip(self, chart_deck, tmp_path):
        out = tmp_path / "nochart.pptx"
        r = run_tool("set_doc_xml", {
            "doc": str(chart_deck), "part": "ppt/charts/chart1.xml",
            "delete": True, "output": str(out),
        })
        assert r["applied"] == 1
        names = [p["part"] for p in run_tool("read_doc_xml", {"doc": str(out)})["parts"]]
        assert not any("charts/chart1.xml" in n for n in names)


class TestLocalMcpServer:
    def test_tool_registry_matches_agent_tools(self):
        from xgen_edit2docs.mcp.local_server import build_local_mcp_server

        mcp = build_local_mcp_server()
        tools = asyncio.run(mcp.list_tools())
        assert sorted(t.name for t in tools) == sorted(TOOL_NAMES)

    def test_deterministic_tool_call_through_mcp(self, deck_path):
        from xgen_edit2docs.mcp.local_server import build_local_mcp_server

        mcp = build_local_mcp_server()
        result = asyncio.run(
            mcp.call_tool("analyze_doc", {"doc": str(deck_path)})
        )
        # FastMCP returns (content_blocks, raw_result)
        raw = result[1] if isinstance(result, tuple) else result
        assert "1" in str(raw) or "page_count" in str(raw)


class TestThemedBuild:
    """DocBuild theme support — a designed deck in ONE deterministic call.

    Closes the gap where keyless agents had no fast path for design-spec
    decks (navy/orange etc.) and fell back to python-pptx scripting or
    20-minute raw-XML authoring."""

    SPEC = {
        "theme": {"bg": "0B1424", "accent": "#EA580C", "rail": True,
                  "page_numbers": True},
        "slides": [
            {"layout": "title", "title": "제목", "subtitle": "발표자"},
            {"layout": "stat", "title": "사용률", "value": "76%",
             "label": "+32%p", "sublabel": "2025"},
            {"layout": "comparison", "title": "두 미래",
             "left": {"heading": "A", "bullets": ["줄1"]},
             "right": {"heading": "B", "bullets": ["줄2"]}},
            {"layout": "quote", "quote": "인용문", "attribution": "출처"},
            {"layout": "content", "title": "행동", "bullets": ["하나", "둘"]},
        ],
    }

    def _build(self, tmp_path):
        out = tmp_path / "themed.pptx"
        run_tool("build_doc", {"spec": self.SPEC, "output": str(out)})
        return Presentation(str(out))

    def test_theme_colors_and_chrome(self, tmp_path):
        prs = self._build(tmp_path)
        assert len(prs.slides) == 5
        # 16:9 widescreen
        assert abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.01
        s1 = list(prs.slides[0].shapes)
        assert str(s1[0].fill.fore_color.rgb) == "0B1424"  # bg
        rail = [sh for sh in s1 if sh.height == prs.slide_height and sh.width < 200000]
        assert rail and str(rail[0].fill.fore_color.rgb) == "EA580C"

    def test_page_numbers_skip_cover(self, tmp_path):
        prs = self._build(tmp_path)
        def texts(i):
            return [sh.text_frame.text for sh in prs.slides[i].shapes
                    if sh.has_text_frame]
        assert not any("/" in t and "01" in t for t in texts(0))  # cover: none
        assert any(t.strip() == "02 / 05" for t in texts(1))

    def test_stat_value_is_accent(self, tmp_path):
        prs = self._build(tmp_path)
        for sh in prs.slides[1].shapes:
            if sh.has_text_frame and "76%" in sh.text_frame.text:
                run = sh.text_frame.paragraphs[0].runs[0]
                assert str(run.font.color.rgb) == "EA580C"
                assert run.font.size.pt == 88
                return
        raise AssertionError("stat value not found")

    def test_comparison_two_panels(self, tmp_path):
        prs = self._build(tmp_path)
        fills = []
        for sh in prs.slides[2].shapes:
            try:
                fills.append(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
        assert fills.count("132339") == 2  # default panel color

    def test_no_theme_keeps_legacy_layouts(self, tmp_path):
        out = tmp_path / "plain.pptx"
        run_tool("build_doc", {"spec": {"slides": [
            {"layout": "title", "title": "T", "subtitle": "S"}]},
            "output": str(out)})
        prs = Presentation(str(out))
        assert prs.slides[0].shapes.title.text == "T"  # placeholder path
